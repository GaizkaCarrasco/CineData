from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colores personalizados
COLOR_PRINCIPAL = RGBColor(31, 78, 121)  # Azul oscuro
COLOR_SECUNDARIO = RGBColor(192, 0, 0)  # Rojo
COLOR_TEXTO = RGBColor(50, 50, 50)
COLOR_BLANCO = RGBColor(255, 255, 255)

def agregar_slide_titulo(prs, titulo, subtitulo=""):
    """Crea un slide con título y subtítulo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRINCIPAL
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = titulo
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_BLANCO
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitulo:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitulo
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = COLOR_SECUNDARIO
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def agregar_slide_contenido(prs, titulo, contenido_items):
    """Crea un slide con título y contenido en viñetas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Encabezado azul
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRINCIPAL
    header_shape.line.color.rgb = COLOR_PRINCIPAL
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = titulo
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_BLANCO
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(contenido_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_TEXTO
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# ==================== SLIDE 1: Objetivo Principal ====================
slide1 = agregar_slide_titulo(
    prs,
    "CineData 🎬",
    "Plataforma de Gestión de Películas con Autenticación JWT"
)

# ==================== SLIDE 2: Objetivo Principal Detallado ====================
objetivo = [
    "🎯 OBJETIVO PRINCIPAL:",
    "",
    "Desarrollar una plataforma web completa y funcional que permita a los",
    "usuarios explorar un catálogo de películas, buscar por criterios",
    "específicos y gestionar sus películas favoritas, con un sistema de",
    "autenticación seguro basado en JWT.",
    "",
    "🎬 VISIÓN DEL PROYECTO:",
    "",
    "Crear un sistema escalable de microservicios que demuestre:",
    "  • Integración de múltiples tecnologías (Python + Node.js)",
    "  • Implementación de seguridad en producción (JWT, bcrypt)",
    "  • Buenas prácticas de arquitectura (separación de responsabilidades)",
    "  • Automatización con Docker (containerización y orquestación)",
    "  • Experiencia de usuario moderna (React SPA)",
    "",
    "🎓 COMPETENCIAS DEMOSTRADAS:",
    "  • Full-Stack Development: Backend + Frontend",
    "  • Microservicios y APIs REST",
    "  • Gestión de datos en múltiples BD (MongoDB + MySQL)",
    "  • DevOps: Docker, Docker Compose"
]
agregar_slide_contenido(prs, "Objetivo Principal", objetivo)

# ==================== SLIDE 3: Modelo de Datos ====================
modelo_datos = [
    "📊 MODELO DE DATOS RELACIONAL Y NO-RELACIONAL:",
    "",
    "📀 MONGODB (NoSQL - Usuarios):",
    "  Colección: users",
    "  └─ _id: ObjectId",
    "  └─ email: string (único)",
    "  └─ password_hash: string (bcrypt)",
    "  └─ is_admin: boolean",
    "  └─ favorites: [movie_ids]",
    "  └─ created_at: timestamp",
    "",
    "🗄️ MYSQL (Relacional - Películas):",
    "  Tabla: movies",
    "  └─ id: INT (PK)",
    "  └─ title: VARCHAR (título)",
    "  └─ genre: VARCHAR (género)",
    "  └─ year: INT (año lanzamiento)",
    "  └─ synopsis: TEXT (sinopsis)",
    "  └─ created_at: TIMESTAMP",
    "",
    "🔑 RELACIÓN ENTRE TABLAS:",
    "  • users.favorites -> almacena IDs de movies",
    "  • Denormalización intencional para optimizar consultas",
    "  • Facilita búsquedas rápidas de favoritos por usuario"
]
agregar_slide_contenido(prs, "Modelo de Datos", modelo_datos)

# ==================== SLIDE 4: Funcionalidades, APIs y Modelo de Datos ====================
funcionalidades = [
    "👤 AUTENTICACIÓN Y USUARIOS:",
    "  • Registro con validación de email",
    "  • Login con generación de JWT",
    "  • Logout con revocación de tokens",
    "  • Contraseñas hasheadas con bcrypt",
    "",
    "🎬 CATÁLOGO DE PELÍCULAS:",
    "  • Listado completo con información detallada",
    "  • Búsqueda en tiempo real por título",
    "  • Filtros por género y año de lanzamiento",
    "  • Vista modal con detalles de película",
    "",
    "⭐ SISTEMA DE FAVORITOS:",
    "  • Guardar películas como favoritas",
    "  • Vista dedicada para favoritos",
    "  • Sincronización en tiempo real",
    "",
    "👨‍💼 PANEL ADMINISTRATIVO:",
    "  • Listar todos los usuarios del sistema",
    "  • Eliminar usuarios (solo admin)",
    "  • Control de acceso restringido",
    "",
    "🔌 APIs REST:",
    "  • UserService: 10+ endpoints (FastAPI en Python)",
    "  • MovieService: 4+ endpoints (Express en Node.js)"
]
agregar_slide_contenido(prs, "Funcionalidades & APIs", funcionalidades)

# ==================== SLIDE 5: Arquitectura y Tecnologías ====================
arquitectura = [
    "🏗️ TIPO DE ARQUITECTURA: MICROSERVICIOS CON API GATEWAY",
    "  • API Gateway centralizado (Express.js en Node.js)",
    "  • 2 servicios backend independientes (desacoplados)",
    "  • 2 bases de datos especializadas (poliglot persistence)",
    "  • Frontend desvinculado de la lógica del servidor",
    "  • Cada servicio puede escalarse independientemente",
    "",
    "🔧 STACK TECNOLÓGICO:",
    "  API Gateway: Express.js en Node.js (puerto 8080)",
    "  Backend 1: FastAPI + Motor (async MongoDB) en Python",
    "  Backend 2: Express.js + MySQL en Node.js",
    "  Frontend: React 19 + Vite + React Router v7",
    "  DevOps: Docker + Docker Compose + Nginx",
    "",
    "🌐 API GATEWAY - PUNTO DE ENTRADA ÚNICO:",
    "  • Autenticación JWT centralizada",
    "  • Proxy inverso inteligente hacia servicios backend",
    "  • Mapeo de rutas: /auth/* → /users/* (users service)",
    "  • Manejo de rutas públicas y protegidas",
    "  • Logging centralizado con Morgan",
    "  • CORS configurado para frontend",
    "",
    "📱 ARQUITECTURA DE FRONTEND:",
    "  • SPA (Single Page Application) con React",
    "  • CSR (Client-Side Rendering)",
    "  • Enrutamiento con React Router v7",
    "  • Comunicación con gateway en puerto 8080",
    "",
    "🔌 FLUJO DE COMUNICACIÓN:",
    "  Frontend (5173) → API Gateway (8080) → Servicios (8000, 3001)"
]
agregar_slide_contenido(prs, "Arquitectura & Tecnologías", arquitectura)

# ==================== SLIDE 6: Lecciones Aprendidas ====================
lecciones = [
    "🎓 LECCIONES APRENDIDAS:",
    "  • Integración eficiente de múltiples lenguajes (Python + Node.js)",
    "  • Orquestación compleja con Docker Compose (6 contenedores)",
    "  • Gestión de dependencias entre servicios",
    "  • Implementación de API Gateway para centralizar autenticación",
    "",
    "⚡ ASPECTOS AVANZADOS IMPLEMENTADOS:",
    "  • API Gateway con proxy inverso inteligente (express-http-proxy)",
    "  • Mapeo de rutas y validación centralizada de JWT",
    "  • Async/await en FastAPI: operaciones no-bloqueantes",
    "  • Motor: driver async para MongoDB con alta concurrencia",
    "  • JWT con tokens con expiración y revocación",
    "  • CORS configurado en gateway y backends para seguridad",
    "  • Rutas protegidas: ProtectedRoute y AdminRoute en React",
    "  • Validación de entrada con Pydantic schemas",
    "  • Hashing seguro de contraseñas con bcrypt",
    "  • Control de acceso basado en roles (RBAC)",
    "",
    "⭐ POR QUÉ MERECE BUENA NOTA:",
    "  ✅ Sistema COMPLETO y FUNCIONAL (Full-Stack + Gateway)",
    "  ✅ Arquitectura PROFESIONAL (3-tier con API Gateway)",
    "  ✅ SEGURIDAD en múltiples niveles (JWT + RBAC + CORS)",
    "  ✅ BUENAS PRÁCTICAS y patrones de diseño actuales",
    "  ✅ ESCALABLE y MANTENIBLE para producción",
    "  ✅ Containerizado con Docker para reproducibilidad"
]
agregar_slide_contenido(prs, "Lecciones & Aspectos Avanzados", lecciones)

# ==================== SLIDE 7: Demo y Conclusiones ====================
demo = [
    "🎬 PASOS DE LA DEMO:",
    "  1. Acceder a http://localhost:5173 (Frontend React)",
    "  2. Registrarse con email y contraseña (JWT generado)",
    "  3. Explorar catálogo completo de películas",
    "  4. Búsqueda en tiempo real por título de película",
    "  5. Filtrar películas por género (Action, Drama, etc.)",
    "  6. Filtrar películas por año de lanzamiento",
    "  7. Hacer clic en película para ver detalles en modal",
    "  8. Guardar/quitar películas de favoritos (⭐)",
    "  9. Ver sección de películas favoritas",
    "  10. Panel Admin (usuario admin): ver y eliminar usuarios",
    "  11. Logout seguro (revocación de token)",
    "",
    "✅ CONCLUSIONES:",
    "  • Plataforma FUNCIONAL y lista para producción",
    "  • Stack MODERNO y ESCALABLE",
    "  • SEGURIDAD implementada en todos los niveles",
    "  • MANTENIBLE y fácil de extender con nuevas features",
    "  • Demuestra DOMINIO de múltiples tecnologías",
    "  • Proyecto que podría comercializarse"
]
agregar_slide_contenido(prs, "Demo & Conclusiones", demo)

# Guardar presentación
prs.save('Presentacion_CineData.pptx')
print("✅ Presentación generada: Presentacion_CineData.pptx")
